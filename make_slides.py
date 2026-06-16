#!/usr/bin/env python3
"""Generate a 2-slide Battleship project summary as a .pptx (opens in Keynote)."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Theme colors
NAVY = RGBColor(0x1A, 0x1A, 0x2E)
CYAN = RGBColor(0x00, 0xD4, 0xFF)
GREY = RGBColor(0x88, 0x92, 0xB0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
RED = RGBColor(0xFF, 0x6B, 0x6B)
PANEL = RGBColor(0x0F, 0x34, 0x60)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
blank = prs.slide_layouts[6]


def bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def textbox(slide, left, top, width, height):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    return tf


def style(run, size, color, bold=False):
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = "Helvetica Neue"


# ---------------- Slide 1 ----------------
s1 = prs.slides.add_slide(blank)
bg(s1, NAVY)

tf = textbox(s1, Inches(0.7), Inches(0.45), Inches(12), Inches(1.1))
p = tf.paragraphs[0]
r = p.add_run(); r.text = "\u2693 Battleship vs. AI"
style(r, 40, CYAN, True)
p2 = tf.add_paragraph()
r = p2.add_run(); r.text = "A browser-based Battleship game you play against an AI opponent"
style(r, 18, GREY)

tf = textbox(s1, Inches(0.7), Inches(1.9), Inches(12), Inches(0.5))
r = tf.paragraphs[0].add_run(); r.text = "What We Built"
style(r, 26, WHITE, True)

points = [
    ("Stack", "Plain HTML, CSS & JavaScript \u2014 no frameworks, no build step, deploys anywhere."),
    ("Gameplay", "Place your 5 ships, then trade shots with the AI until one fleet is sunk."),
    ("Smart AI", "Fires randomly until it hits, then enters \u201ctarget mode\u201d to hunt adjacent cells."),
    ("Drag-to-place", "Drag sideways for horizontal, up/down for vertical \u2014 with live preview & undo."),
    ("Polished UX", "Name entry, turn indicator, hit/miss/sunk messages, start validation popup."),
    ("Responsive", "Fits the screen without scrolling; works on desktop & touch."),
    ("Shipped", "Live on the web + full source in a public GitHub repo."),
]
tf = textbox(s1, Inches(0.7), Inches(2.5), Inches(12), Inches(4.6))
first = True
for title, desc in points:
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    first = False
    p.space_after = Pt(8)
    r = p.add_run(); r.text = f"\u2022  {title}:  "
    style(r, 17, CYAN, True)
    r = p.add_run(); r.text = desc
    style(r, 17, WHITE)


# ---------------- Slide 2 ----------------
s2 = prs.slides.add_slide(blank)
bg(s2, NAVY)

tf = textbox(s2, Inches(0.7), Inches(0.45), Inches(12), Inches(0.8))
r = tf.paragraphs[0].add_run(); r.text = "Bugs Found & How We Fixed Them"
style(r, 30, CYAN, True)

rows = [
    ("Bug", "Root Cause", "Fix"),
    ("AI turn crashed the game",
     "`cell` declared inside a loop but used outside it \u2192 ReferenceError each AI move",
     "Re-declared cell in the correct scope"),
    ("AI could freeze / overflow",
     "AI re-picked hit cells via infinite recursion",
     "Replaced recursion with a retry while-loop"),
    ("Ship counts wrong across games",
     "playerShips array never reset on a new game",
     "Reset all placement state in initSetup()"),
    ("Hit counts overflowed",
     "A sunk ship kept incrementing its total",
     "Guard: stop counting at hits >= length"),
    ("Buttons \u201cdidn\u2019t work\u201d",
     "Listeners ran before the DOM existed",
     "Centralized wiring in init() on DOMContentLoaded"),
    ("Edits never appeared",
     "Dev server served cached 304s (stale files)",
     "No-cache dev server + versioned assets"),
    ("Had to scroll to play",
     "Fixed sizing overflowed shorter screens",
     "Sized board from leftover viewport height"),
]

left, top = Inches(0.6), Inches(1.45)
width, height = Inches(12.1), Inches(5.0)
table = s2.shapes.add_table(len(rows), 3, left, top, width, height).table
table.columns[0].width = Inches(3.3)
table.columns[1].width = Inches(4.9)
table.columns[2].width = Inches(3.9)

for ci in range(3):
    cell = table.cell(0, ci)
    cell.fill.solid(); cell.fill.fore_color.rgb = CYAN
    para = cell.text_frame.paragraphs[0]
    run = para.add_run(); run.text = rows[0][ci]
    style(run, 15, NAVY, True)

for ri in range(1, len(rows)):
    for ci in range(3):
        cell = table.cell(ri, ci)
        cell.fill.solid()
        cell.fill.fore_color.rgb = PANEL if ri % 2 else NAVY
        para = cell.text_frame.paragraphs[0]
        run = para.add_run(); run.text = rows[ri][ci]
        style(run, 12.5, WHITE if ci != 0 else CYAN, ci == 0)

tf = textbox(s2, Inches(0.6), Inches(6.7), Inches(12.1), Inches(0.6))
r = tf.paragraphs[0].add_run()
r.text = ("Process: found issues via code review + live testing, fixed the root cause "
          "(not symptoms) each time, and verified in-browser after every change.")
style(r, 13, GREY)

out = "Battleship_Summary.pptx"
prs.save(out)
print("Saved", out)
